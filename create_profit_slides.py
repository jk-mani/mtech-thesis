"""Create PowerPoint slides for GT0 DQN experiments and Profit comparison."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle indentation
        if line.startswith("  "):
            p.text = "  • " + line.strip()
            p.level = 1
            p.font.size = Pt(18)
        else:
            p.text = "• " + line
            p.font.size = Pt(20)
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.2), Inches(1.3), Inches(9.6), Inches(0.4 * num_rows)).table
    
    # Set column widths based on number of columns
    if num_cols <= 4:
        table.columns[0].width = Inches(2.5)
        for i in range(1, num_cols):
            table.columns[i].width = Inches(2.3)
    else:
        # Narrower columns for tables with many columns
        table.columns[0].width = Inches(1.8)
        for i in range(1, num_cols):
            table.columns[i].width = Inches(1.4)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_image_slide(prs, title, image_path, caption=None):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Image - positioned with proper spacing for title and caption
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.3), width=Inches(9), height=Inches(4.8))
    else:
        # Placeholder text if image not found
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image: {Path(image_path).name}]"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
    
    # Caption - positioned at bottom with more clearance
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Profit-Based DQN for Bike Rebalancing",
        "Comparing Profit vs Lost-Demand Objectives"
    )
    
    # Slide 2: Problem Statement
    add_content_slide(prs, "Why Profit-Based Rewards?", [
        "Current Approach: Minimize lost demand",
        "  Ignores operational costs (truck travel, labor)",
        "  May lead to excessive, unprofitable rebalancing",
        "",
        "Proposed Approach: Maximize profit",
        "  Profit = Revenue - Truck Cost - Lost Demand Penalty",
        "  Agent learns cost-efficient rebalancing strategies",
        "",
        "Goal: Show profit-DQN achieves better ROI"
    ])
    
    # Slide 3: Economic Model
    add_content_slide(prs, "Economic Model", [
        "Trip Revenue (Distance-Based Pricing):",
        "  Revenue = $1.00 (base) + $0.75 × trip_km",
        "",
        "Truck Operating Cost:",
        "  Cost = $1.00 × truck_distance_km",
        "",
        "Lost Demand Penalty:",
        "  Penalty = $5.00 per lost trip",
        "",
        "Profit = Σ(Trip Revenue) - Truck Cost - Lost Penalty"
    ])
    
    # Slide 4: Experimental Setup
    add_content_slide(prs, "Experimental Setup", [
        "Environment: GT0 Toy Model",
        "  10 stations, 2 vehicles, 15-bike capacity",
        "",
        "Training: 50,000 timesteps each",
        "  Profit-DQN (PReLU): Trained with profit reward",
        "  Lost-Demand-DQN (ELU): Trained with -lost_demand reward",
        "",
        "Evaluation: 50 test episodes (unseen data)",
        "  Both models evaluated on same test scenarios"
    ])
    
    # Slide 5: Policy Comparison Results (with profit)
    add_table_slide(
        prs,
        "Policy Comparison Results (50 Test Episodes)",
        ["Metric", "Lost-Demand DQN", "Profit DQN", "Difference"],
        [
            ["Avg Profit", "$113.95", "$128.01", "+$14.06 (better)"],
            ["Avg Revenue", "$194.87", "$198.20", "+$3.33"],
            ["Avg Truck Cost", "$57.22", "$49.59", "-$7.63 (less)"],
            ["Avg Lost Penalty", "$23.70", "$20.60", "-$3.10"],
            ["Truck Distance", "57.22 km", "49.59 km", "-7.63 km"],
            ["Lost Demand Rate", "6.60%", "5.96%", "-0.64%"]
        ]
    )
    
    # Slide 6: Profit Breakdown
    base_dir = Path(__file__).parent
    add_image_slide(
        prs,
        "Profit Breakdown: Revenue vs Costs",
        base_dir / "results_policy_comparison" / "profit_breakdown.png",
        "Profit DQN reduces truck costs while maintaining revenue"
    )
    
    # Slide 7: Net Profit Comparison
    add_image_slide(
        prs,
        "Net Profit Comparison",
        base_dir / "results_policy_comparison" / "profit_comparison.png",
        "Profit DQN achieves $14.06 higher average profit per episode"
    )
    
    # Slide 8: Service Quality vs Economics
    add_image_slide(
        prs,
        "Service Quality vs Economic Performance",
        base_dir / "results_policy_comparison" / "policy_summary.png",
        "Profit DQN balances service quality with economic efficiency"
    )
    
    # Slide 9: Economic Sensitivity - Introduction
    add_content_slide(prs, "Economic Parameter Sensitivity", [
        "Research Question:",
        "  How do economic parameters affect learned policies?",
        "",
        "Parameters varied:",
        "  Trip revenue: $0.50-$2.00 base + $0.50-$1.00/km",
        "  Truck cost: $0.50, $1.00, $2.00 per km",
        "  Lost demand penalty: $2.00, $5.00, $10.00 per trip",
        "",
        "Each config trained 50,000 timesteps, evaluated 50 episodes"
    ])
    
    # Slide 10: Economic Sensitivity - Lost Demand Plot
    add_image_slide(
        prs,
        "Economic Sensitivity: Lost Demand Rate",
        base_dir / "results_economic_sensitivity" / "economic_lost_demand.png",
        "Lower truck cost and baseline penalty achieve best service quality"
    )
    
    # Slide 11: Combined Economic & Revenue Sensitivity Results
    add_table_slide(
        prs,
        "Economic Parameter Sensitivity: Full Results",
        ["Configuration", "Revenue", "Cost/km", "Penalty", "Profit", "Lost %"],
        [
            ["baseline", "$1+$0.75/km", "$1.00", "$5.00", "$128.01", "5.96%"],
            ["high_revenue", "$2+$1/km", "$1.00", "$5.00", "$243.00", "5.92%"],
            ["low_revenue", "$0.5+$0.5/km", "$1.00", "$5.00", "$48.33", "7.49%"],
            ["low_cost", "$1+$0.75/km", "$0.50", "$5.00", "$149.09", "5.86%"],
            ["high_cost", "$1+$0.75/km", "$2.00", "$5.00", "$80.93", "5.87%"],
            ["low_penalty", "$1+$0.75/km", "$1.00", "$2.00", "$135.35", "7.43%"],
            ["high_penalty", "$1+$0.75/km", "$1.00", "$10.00", "$93.83", "8.22%"],
            ["cheap_aggressive", "$1+$0.75/km", "$0.50", "$10.00", "$124.60", "6.13%"],
            ["expensive_cons.", "$1+$0.75/km", "$2.00", "$2.00", "$97.88", "7.53%"]
        ]
    )
    
    # Slide 12: Key Findings
    add_content_slide(prs, "Key Findings", [
        "Policy Comparison:",
        "  Profit DQN achieves $14.06 higher profit per episode",
        "  Reduces truck costs by 13%, improves service by 0.64%",
        "",
        "Economic Parameter Impact (ranked by profit):",
        "  1. Revenue pricing: Strongest impact ($48 → $243)",
        "  2. Truck cost: Moderate ($81 → $149 at low cost)",
        "  3. Lost penalty: Inverse effect (high penalty hurts)",
        "",
        "Best config: high_revenue ($243), Worst: low_revenue ($48)"
    ])
    
    # Slide 13: Conclusion
    add_content_slide(prs, "Conclusions", [
        "Profit-Based DQN Benefits:",
        "  Better service: 5.96% vs 6.60% lost demand",
        "  Cost-aware decisions: Balances service vs efficiency",
        "",
        "Economic Parameter Insights:",
        "  Moderate penalty ($5) is optimal",
        "  Extreme penalties hurt convergence",
        "  Truck cost less impactful than expected",
        "",
        "Recommendations:",
        "  Use profit-based objectives for real deployments",
        "  Tune lost demand penalty carefully"
    ])
    
    # Save
    output_path = Path(__file__).parent / "Profit_DQN_Comparison.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
